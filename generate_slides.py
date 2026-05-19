"""
Génère le PowerPoint de présentation du projet (~18 slides).
Lancer : python generate_slides.py
Sortie : Presentation_Projet4.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Palette ─────────────────────────────────────────────────
PRIMARY = RGBColor(0x1a, 0x3a, 0x6e)
SECONDARY = RGBColor(0x2c, 0x5a, 0xa0)
ACCENT = RGBColor(0xd9, 0x77, 0x06)
SUCCESS = RGBColor(0x15, 0x80, 0x3d)
DANGER = RGBColor(0xb9, 0x1c, 0x1c)
LIGHT_BG = RGBColor(0xf1, 0xf5, 0xf9)
DARK = RGBColor(0x0f, 0x17, 0x2a)
GRAY = RGBColor(0x64, 0x74, 0x8b)
WHITE = RGBColor(0xff, 0xff, 0xff)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Layout blank
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=LIGHT_BG):
    """Ajoute un fond coloré."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                     prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_title_bar(slide, title, subtitle=""):
    """Barre de titre en haut."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                   prs.slide_width, Inches(1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    bar.shadow.inherit = False

    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(14)
        r2.font.color.rgb = LIGHT_BG

def add_text(slide, text, left, top, width, height,
             size=14, color=DARK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, bg=None):
    """Ajoute une zone de texte."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    if bg:
        tb.fill.solid()
        tb.fill.fore_color.rgb = bg
        tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    # Handle multi-line text
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
    return tb

def add_bullets(slide, items, left, top, width, height, size=14,
                color=DARK, bullet_color=ACCENT):
    """Ajoute une liste à puces."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        r1 = p.add_run()
        r1.text = "▸ "
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color

def add_footer(slide, page_num, total):
    """Footer avec page number."""
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(7.1),
                                    Inches(12.7), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"Projet 4 — SMA Agentic AI · 4AISDR 2026"
    r.font.size = Pt(9)
    r.font.color.rgb = GRAY

    p2 = tf.paragraphs[0]
    r3 = p2.add_run()
    r3.text = f"   ·   {page_num}/{total}"
    r3.font.size = Pt(9)
    r3.font.color.rgb = GRAY


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Titre
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, PRIMARY)

# Logo / emoji
add_text(slide, "🎓", Inches(5.5), Inches(0.8), Inches(2.5), Inches(1.5),
         size=80, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "Learning Path Architect",
         Inches(1), Inches(2.5), Inches(11.5), Inches(0.8),
         size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, "Système Multi-Agents pour parcours d'apprentissage adaptatifs",
         Inches(1), Inches(3.4), Inches(11.5), Inches(0.6),
         size=20, color=LIGHT_BG, italic=True, align=PP_ALIGN.CENTER)

# Séparateur
sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(5), Inches(4.3), Inches(3.3), Inches(0.05))
sep.fill.solid()
sep.fill.fore_color.rgb = ACCENT
sep.line.fill.background()

add_text(slide, "Projet 4   ·   Module 4AISDR   ·   Mai 2026",
         Inches(1), Inches(4.6), Inches(11.5), Inches(0.5),
         size=18, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide,
         "Équipe : Othmane Dhilou  ·  Dodo  ·  Mohamed Yassir\n"
         "Encadrants : Pr. Hasnââ CHAABI  ·  Pr. Nadia IDRISSI Zouggari",
         Inches(1), Inches(5.5), Inches(11.5), Inches(1.2),
         size=14, color=LIGHT_BG, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Sommaire
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_title_bar(slide, "Plan de la présentation", "Comment nous allons couvrir le projet en 20 minutes")

bullets = [
    "Contexte & Problématique",
    "Architecture du Système Multi-Agents",
    "Démonstration live (3 profils)",
    "Choix techniques justifiés",
    "Conformité au cours (Chapitres 1 à 7)",
    "Limitations & perspectives",
    "Questions / Réponses"
]
add_bullets(slide, bullets, Inches(2), Inches(1.6), Inches(9), Inches(5),
            size=20, bullet_color=ACCENT)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Problématique
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_title_bar(slide, "Contexte & Problématique", "Pourquoi un SMA pour l'apprentissage adaptatif ?")

add_text(slide, "Le problème",
         Inches(0.7), Inches(1.4), Inches(6), Inches(0.5),
         size=22, color=DANGER, bold=True)

add_bullets(slide, [
    "Diversité croissante des profils d'étudiants",
    "L'approche unique pour tous montre ses limites",
    "Difficulté d'identifier les vraies lacunes",
    "Surcharge des enseignants pour personnaliser"
], Inches(0.7), Inches(2), Inches(6), Inches(3), size=15, bullet_color=DANGER)

add_text(slide, "La solution",
         Inches(7.2), Inches(1.4), Inches(5.5), Inches(0.5),
         size=22, color=SUCCESS, bold=True)

add_bullets(slide, [
    "5 agents spécialisés qui collaborent",
    "Test adaptatif (IRT simplifié)",
    "Parcours personnalisé respectant pré-requis",
    "RAG vectoriel pour ressources pertinentes",
    "Répétition espacée SM-2 (algorithme Anki)"
], Inches(7.2), Inches(2), Inches(5.5), Inches(3), size=15, bullet_color=SUCCESS)

add_text(slide, "💡 Un système qui s'adapte à chaque apprenant — automatiquement",
         Inches(1), Inches(6), Inches(11.3), Inches(0.6),
         size=18, color=PRIMARY, bold=True, italic=True, align=PP_ALIGN.CENTER,
         bg=LIGHT_BG)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture vue d'ensemble
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_title_bar(slide, "Architecture multi-agents", "5 agents spécialisés orchestrés par LangGraph")

# 5 cartes d'agents
agents = [
    ("🩺", "Diagnosticien", "Test adaptatif\n+ identification\ndes lacunes", PRIMARY),
    ("🗺️", "Planificateur", "Parcours ordonné\n+ Chain-of-Thought\n(Mind Layer)", SECONDARY),
    ("📚", "Pédagogue", "RAG vectoriel\nChroma +\nembeddings", ACCENT),
    ("🔁", "Coach", "Répétition\nespacée SM-2\n(Anki)", SUCCESS),
    ("💾", "Tracker", "Mémoire\nlong-terme\nJSON", DARK),
]
for i, (emoji, nom, desc, color) in enumerate(agents):
    left = Inches(0.5 + i * 2.5)
    # Carte
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, Inches(1.5), Inches(2.3), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = color
    card.line.width = Pt(2)
    card.shadow.inherit = False

    # Header coloré
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      left, Inches(1.5), Inches(2.3), Inches(0.6))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    header.shadow.inherit = False

    # Emoji
    add_text(slide, emoji, left, Inches(2.2), Inches(2.3), Inches(0.8),
             size=40, align=PP_ALIGN.CENTER)
    # Nom
    add_text(slide, nom, left, Inches(1.55), Inches(2.3), Inches(0.5),
             size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Description
    add_text(slide, desc, left, Inches(3.2), Inches(2.3), Inches(1.5),
             size=11, color=DARK, align=PP_ALIGN.CENTER)

add_text(slide,
         "Communication inter-agents via protocole A2A inspiré FIPA-ACL\n"
         "Orchestration par LangGraph avec State partagé + Conditional Edges",
         Inches(0.5), Inches(5.5), Inches(12.3), Inches(1),
         size=15, color=PRIMARY, bold=True, italic=True, align=PP_ALIGN.CENTER,
         bg=LIGHT_BG)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Le graphe LangGraph
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_title_bar(slide, "Orchestration LangGraph", "State + Nodes + Edges + Conditional Edge")

# ASCII representation
add_text(slide,
"""                    ┌─────────┐
                    │  START  │
                    └────┬────┘
                         ▼
                ┌────────────────┐
                │  Diagnosticien │
                └────────┬───────┘
                         │
                  ╔══════╧══════╗
                  ║ Conditional ║   ← cœur LangGraph !
                  ╚══════╤══════╝
                ≥80%    │    <80%
              ┌────────┘      └────────┐
              ▼                        ▼
        ┌──────────┐         ┌────────────────┐
        │  Coach   │         │ Planificateur  │
        │ (direct) │         └────────┬───────┘
        └────┬─────┘                  ▼
             │              ┌──────────────────┐
             │              │    Pédagogue     │
             │              └──────────┬───────┘
             │                         ▼
             │                  ┌──────────┐
             │                  │  Coach   │
             │                  └────┬─────┘
             └──────────────────────┘
                              ▼
                        ┌──────────┐
                        │ Tracker  │
                        └─────┬────┘
                              ▼
                          ┌───────┐
                          │  END  │
                          └───────┘""",
         Inches(0.5), Inches(1.4), Inches(8), Inches(5.5),
         size=11, color=DARK, align=PP_ALIGN.LEFT)

add_text(slide,
         "Conditional Edge\n(valeur ajoutée majeure de LangGraph vs LangChain)",
         Inches(9), Inches(1.4), Inches(4), Inches(1),
         size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER,
         bg=LIGHT_BG)

add_text(slide,
"""def route_after_diagnostic(state):
    if state['diagnostic']['pourcentage'] >= 80:
        return 'coach'      # skip planif !
    return 'planificateur'""",
         Inches(9), Inches(2.6), Inches(4), Inches(2),
         size=10, color=DARK)

add_bullets(slide, [
    "Si score ≥ 80% : skip planif + pédagogue",
    "Sinon : parcours complet",
    "Économie démontrable en démo",
    "Conforme cours Ch.5 LangGraph"
], Inches(9), Inches(4.5), Inches(4), Inches(2.5),
   size=12, bullet_color=SUCCESS)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Communication A2A
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_title_bar(slide, "Communication inter-agents (A2A)", "Protocole inspiré FIPA-ACL")

add_text(slide, "Format de message standardisé",
         Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
         size=18, color=PRIMARY, bold=True)

add_text(slide,
"""{
  "msg_id": "uuid-1234",
  "timestamp": "10:23:45",
  "from": "Orchestrateur",
  "to": "Diagnosticien",
  "performative": "REQUEST",
  "content": {
    "action": "run_diagnostic",
    "module": "MSA"
  },
  "conversation_id": "session-001"
}""",
         Inches(0.5), Inches(1.9), Inches(6), Inches(4.5),
         size=12, color=DARK, bg=LIGHT_BG)

add_text(slide, "Performatives FIPA-ACL utilisées",
         Inches(7), Inches(1.3), Inches(6), Inches(0.5),
         size=18, color=PRIMARY, bold=True)

perfs = [
    ("REQUEST", "demande une action", PRIMARY),
    ("INFORM", "transmet une information", SUCCESS),
    ("QUERY", "demande une information", ACCENT),
    ("CONFIRM", "confirme un fait", SECONDARY),
    ("FAILURE", "signale un échec", DANGER),
]
for i, (perf, desc, color) in enumerate(perfs):
    add_text(slide, f"{perf}", Inches(7), Inches(2 + i*0.6),
             Inches(2), Inches(0.5),
             size=14, color=color, bold=True)
    add_text(slide, desc, Inches(9), Inches(2 + i*0.6),
             Inches(4), Inches(0.5),
             size=12, color=DARK)

add_text(slide,
         "💡 Visible en direct dans la timeline Streamlit pendant la démo",
         Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
         size=14, color=ACCENT, italic=True, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — RAG vectoriel
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_title_bar(slide, "RAG vectoriel", "Chroma + Embeddings multilingues (Ch. 2 du cours)")

add_text(slide, "Pipeline d'indexation",
         Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
         size=18, color=PRIMARY, bold=True)

pipeline = [
    "1. Charger ressources.json + textes .txt",
    "2. Chunking récursif (¶ → ligne → phrase → mot)",
    "3. Embedding via multilingual-MiniLM (384 dim)",
    "4. Indexation Chroma (persistance SQLite)",
    "5. À la query : similarité cosine + top-K"
]
add_bullets(slide, pipeline, Inches(0.5), Inches(1.9),
             Inches(6), Inches(3.5), size=14, bullet_color=ACCENT)

add_text(slide, "Justification des choix",
         Inches(7), Inches(1.3), Inches(6), Inches(0.5),
         size=18, color=PRIMARY, bold=True)

justifs = [
    "Chroma : persistance auto + métadonnées",
    "multilingual-MiniLM : 384 dim, FR + EN",
    "Gratuit, local (pas d'API requise)",
    "RecursiveCharacterTextSplitter (cours Ch.2)",
    "Cosine similarity (standard du domaine)"
]
add_bullets(slide, justifs, Inches(7), Inches(1.9),
             Inches(6), Inches(3.5), size=14, bullet_color=SUCCESS)

add_text(slide,
         "Exemple : query='K-Means partitionnement' → score 0.634 sur 'K-Means Clustering'",
         Inches(0.5), Inches(6), Inches(12.3), Inches(0.6),
         size=13, color=PRIMARY, italic=True, align=PP_ALIGN.CENTER, bg=LIGHT_BG)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — MCP
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_title_bar(slide, "MCP (Model Context Protocol)", "Standard ouvert Anthropic — l'USB-C du LLM")

add_text(slide,
         "« MCP est au LLM ce que l'USB-C est au téléphone »",
         Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.6),
         size=18, color=ACCENT, italic=True, bold=True, align=PP_ALIGN.CENTER,
         bg=LIGHT_BG)

add_text(slide, "3 outils @mcp.tool() exposés :",
         Inches(0.5), Inches(2.2), Inches(6), Inches(0.5),
         size=16, color=PRIMARY, bold=True)

tools = [
    ("search_ressources_rag(query, n)", "→ Recherche RAG via Chroma"),
    ("compute_sm2_schedule(notion, q)", "→ Algorithme de répétition espacée"),
    ("update_profile_ltm(etudiant)", "→ Mémoire long-terme JSON"),
]
for i, (tool, desc) in enumerate(tools):
    add_text(slide, tool, Inches(0.5), Inches(2.8 + i*0.7),
             Inches(5), Inches(0.5), size=13, color=DARK, bold=True)
    add_text(slide, desc, Inches(5.5), Inches(2.8 + i*0.7),
             Inches(5), Inches(0.5), size=12, color=GRAY)

add_text(slide, "Avantages dans notre architecture",
         Inches(7.5), Inches(2.2), Inches(5.5), Inches(0.5),
         size=16, color=PRIMARY, bold=True)

add_bullets(slide, [
    "Isolation processus (sous-process stdio)",
    "Démarrage embedding ne bloque pas l'UI",
    "Réutilisable par d'autres projets",
    "Conforme au standard de l'industrie",
    "Sécurité (variables d'env, pas de secrets en dur)"
], Inches(7.5), Inches(2.8), Inches(5.5), Inches(3),
   size=12, bullet_color=SUCCESS)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Mind Layer (LLM + CoT)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_title_bar(slide, "Mind Layer — LLM + Chain-of-Thought", "Conforme Ch.1 (Prompt Engineering) + Ch.3 (Agentic AI)")

add_text(slide, "Le Planificateur a un cerveau (Mind)",
         Inches(0.5), Inches(1.3), Inches(12), Inches(0.6),
         size=20, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, "Architecture",
         Inches(0.5), Inches(2.1), Inches(6), Inches(0.5),
         size=16, color=ACCENT, bold=True)

add_bullets(slide, [
    "Détection auto du LLM disponible :",
    "   OpenRouter / OpenAI / Gemini",
    "Fallback déterministe (templates CoT)",
    "Aucun secret en dur (variables d'env)",
    "@wrap_model_call compatible (Ch.4)"
], Inches(0.5), Inches(2.6), Inches(6), Inches(3), size=13, bullet_color=ACCENT)

add_text(slide, "Exemple de raisonnement CoT généré",
         Inches(7), Inches(2.1), Inches(6), Inches(0.5),
         size=16, color=ACCENT, bold=True)

add_text(slide,
"""Étape 1 — Analyse : Score 50%
avec 2 lacunes sur notions avancées.

Étape 2 — Identification des
prérequis : les notions maîtrisées
servent de fondation.

Étape 3 — Construction du parcours :
2 étapes ordonnées par complexité.

Étape 4 — SM-2 : intervalles
progressifs (1j → 6j → 15j)""",
         Inches(7), Inches(2.6), Inches(6), Inches(4),
         size=11, color=DARK, bg=LIGHT_BG)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 — Démo profil 1 : Débutant
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_title_bar(slide, "Démo 1 — Amina (débutante)", "Score 0% → parcours complet")

add_text(slide, "Profil",
         Inches(0.5), Inches(1.3), Inches(4), Inches(0.5),
         size=18, color=PRIMARY, bold=True)
add_bullets(slide, [
    "Amina, débutante",
    "Module : MSA",
    "Score initial : 0%",
    "4 lacunes détectées"
], Inches(0.5), Inches(1.9), Inches(4), Inches(2.5),
   size=13, bullet_color=DANGER)

add_text(slide, "Chemin parcouru",
         Inches(5), Inches(1.3), Inches(4), Inches(0.5),
         size=18, color=PRIMARY, bold=True)
add_bullets(slide, [
    "Diagnosticien (4 questions)",
    "Planificateur (parcours 4 étapes)",
    "Pédagogue (5 ressources RAG)",
    "Coach (SM-2 → demain)",
    "Tracker (sauvé profil)"
], Inches(5), Inches(1.9), Inches(4), Inches(3.5),
   size=13, bullet_color=PRIMARY)

add_text(slide, "Résultats",
         Inches(9.5), Inches(1.3), Inches(3.3), Inches(0.5),
         size=18, color=PRIMARY, bold=True)
add_bullets(slide, [
    "📊 4 étapes parcours",
    "📚 5 ressources",
    "📨 11 messages A2A",
    "📅 Révision demain"
], Inches(9.5), Inches(1.9), Inches(3.3), Inches(3),
   size=13, bullet_color=SUCCESS)

add_text(slide,
         "💡 Démo live — montrer la timeline A2A et le raisonnement CoT",
         Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6),
         size=14, color=ACCENT, italic=True, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 — Démo profil 3 : Avancé (Conditional Edge !)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_title_bar(slide, "Démo 3 — Sarah (avancée)",
              "Score 100% → Conditional Edge déclenchée (saut du parcours)")

add_text(slide, "Profil",
         Inches(0.5), Inches(1.3), Inches(4), Inches(0.5),
         size=18, color=PRIMARY, bold=True)
add_bullets(slide, [
    "Sarah, étudiante avancée",
    "Module : MSA",
    "Score initial : 100%",
    "Aucune lacune"
], Inches(0.5), Inches(1.9), Inches(4), Inches(2.5),
   size=13, bullet_color=SUCCESS)

add_text(slide, "Chemin parcouru",
         Inches(5), Inches(1.3), Inches(4), Inches(0.5),
         size=18, color=PRIMARY, bold=True)
add_bullets(slide, [
    "Diagnosticien (4 questions)",
    "❌ Planificateur SKIPPED",
    "❌ Pédagogue SKIPPED",
    "✅ Coach direct (SM-2)",
    "✅ Tracker"
], Inches(5), Inches(1.9), Inches(4), Inches(3.5),
   size=13, bullet_color=ACCENT)

add_text(slide, "Économie",
         Inches(9.5), Inches(1.3), Inches(3.3), Inches(0.5),
         size=18, color=PRIMARY, bold=True)
add_bullets(slide, [
    "📊 0 étape parcours",
    "📚 0 ressource",
    "📨 7 msg A2A (vs 11)",
    "⏱ -36% de calculs",
    "🎯 Conditional !"
], Inches(9.5), Inches(1.9), Inches(3.3), Inches(3),
   size=13, bullet_color=SUCCESS)

add_text(slide,
         "🏆 C'est ICI qu'on prouve la valeur ajoutée de LangGraph",
         Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6),
         size=14, color=DANGER, italic=True, bold=True, align=PP_ALIGN.CENTER,
         bg=LIGHT_BG)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 — Stack technique récap
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_title_bar(slide, "Stack technique", "Une technologie par exigence du cours")

stack = [
    ("Orchestration", "LangGraph", "Ch.5 — State + Nodes + Conditional"),
    ("LLM (Mind)", "OpenRouter / fallback", "Ch.1 — Chain-of-Thought"),
    ("Outils (Body)", "MCP + FastMCP", "Ch.6 — @mcp.tool"),
    ("RAG", "Chroma + multilingual-MiniLM", "Ch.2 — chunking + embeddings"),
    ("Mémoire CT", "InMemorySaver", "Ch.3 — checkpointing"),
    ("Mémoire LT", "JSON files", "Ch.3 — persistence"),
    ("Répétition", "SuperMemo-2 (Anki)", "Sujet du projet"),
    ("Comm. agents", "A2A FIPA-ACL", "Ch.7 — SMA & A2A"),
    ("UI", "Streamlit", "Démo visuelle"),
    ("Tests", "pytest (15 tests)", "Qualité code"),
]
for i, (couche, tech, justif) in enumerate(stack):
    y = Inches(1.3 + i * 0.5)
    add_text(slide, couche, Inches(0.5), y, Inches(3), Inches(0.4),
             size=13, color=PRIMARY, bold=True)
    add_text(slide, tech, Inches(3.7), y, Inches(4), Inches(0.4),
             size=13, color=DARK)
    add_text(slide, justif, Inches(8), y, Inches(5), Inches(0.4),
             size=12, color=GRAY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 13 — Conformité au cours
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_title_bar(slide, "Conformité au cours", "Mapping projet ↔ chapitres")

mapping = [
    ("Ch.1 — Prompt Engineering (CoT, ReAct)", "✅", "mind_layer.py"),
    ("Ch.2 — RAG (Chunking + Embeddings)", "✅", "rag_engine.py (Chroma)"),
    ("Ch.3 — Agentic AI (Body/Mind/Memory)", "✅", "orchestrator.py"),
    ("Ch.4 — LangChain + Middlewares", "✅", "wrap_model_call ready"),
    ("Ch.5 — LangGraph (State/Nodes/Edges)", "✅", "orchestrator.py"),
    ("Ch.6 — MCP (FastMCP + @mcp.tool)", "✅", "learning_tools_server.py"),
    ("Ch.7 — SMA & A2A (FIPA-ACL)", "✅", "a2a_protocol.py"),
]
for i, (chap, status, file) in enumerate(mapping):
    y = Inches(1.4 + i * 0.65)
    # Status box
    status_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          Inches(0.5), y, Inches(0.6), Inches(0.5))
    status_box.fill.solid()
    status_box.fill.fore_color.rgb = SUCCESS
    status_box.line.fill.background()
    status_box.shadow.inherit = False
    add_text(slide, status, Inches(0.5), y, Inches(0.6), Inches(0.5),
             size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, chap, Inches(1.3), y, Inches(7), Inches(0.5),
             size=14, color=DARK, bold=True)
    add_text(slide, file, Inches(8.5), y, Inches(4.5), Inches(0.5),
             size=12, color=GRAY, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 14 — Pourquoi ces choix
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_title_bar(slide, "Pourquoi ces choix ?", "Anticipation des questions du jury")

choices = [
    ("Pourquoi LangGraph ?",
     "Logique non linéaire avec conditional edges (impossible en LangChain pur)"),
    ("Pourquoi Chroma et pas FAISS ?",
     "Métadonnées + persistance + 20 ressources (FAISS = surdesign)"),
    ("Pourquoi MCP en plus ?",
     "Chapitre majeur du cours + isolation + réutilisabilité"),
    ("Pourquoi SM-2 et pas règles simples ?",
     "Algorithme de référence (Anki), exigé par le sujet"),
    ("Pourquoi un LLM fallback ?",
     "Démo garantie même sans Internet (jour J)"),
    ("Pourquoi A2A FIPA-ACL ?",
     "Standard SMA + visibilité au jury via timeline"),
]
for i, (q, a) in enumerate(choices):
    y = Inches(1.3 + i * 0.85)
    add_text(slide, q, Inches(0.5), y, Inches(5), Inches(0.5),
             size=14, color=ACCENT, bold=True)
    add_text(slide, a, Inches(5.5), y, Inches(7.3), Inches(0.7),
             size=12, color=DARK)

# ════════════════════════════════════════════════════════════════
# SLIDE 15 — Tests
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_title_bar(slide, "Tests unitaires", "15 tests pytest — qualité code")

add_text(slide, "Lancement",
         Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
         size=18, color=PRIMARY, bold=True)
add_text(slide, "$ pytest tests/ -v",
         Inches(0.5), Inches(1.9), Inches(6), Inches(0.6),
         size=14, color=WHITE, bold=True, bg=DARK)
add_text(slide, "→ 15 passed in ~1s ✅",
         Inches(0.5), Inches(2.7), Inches(6), Inches(0.5),
         size=14, color=SUCCESS, bold=True)

add_text(slide, "Couverture",
         Inches(0.5), Inches(3.4), Inches(6), Inches(0.5),
         size=18, color=PRIMARY, bold=True)

add_bullets(slide, [
    "TestConditionalEdge (5 tests)",
    "TestA2AProtocol (4 tests)",
    "TestSessionsIntegration (5 tests)",
    "TestNodesIsolation (1 test)"
], Inches(0.5), Inches(4), Inches(6), Inches(3),
   size=13, bullet_color=SUCCESS)

add_text(slide, "Exemple de test critique",
         Inches(7), Inches(1.3), Inches(6), Inches(0.5),
         size=18, color=PRIMARY, bold=True)

add_text(slide,
"""def test_avance_route_vers_coach():
    \"\"\"Score >= 80% doit
    router vers coach
    (skip planificateur).\"\"\"
    state = {
        'diagnostic': {
            'pourcentage': 85
        }
    }
    assert route_after_diagnostic(state) == 'coach'""",
         Inches(7), Inches(1.9), Inches(6), Inches(5),
         size=11, color=DARK, bg=LIGHT_BG)

# ════════════════════════════════════════════════════════════════
# SLIDE 16 — Limitations & perspectives
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_title_bar(slide, "Limitations & Perspectives", "Honnêteté intellectuelle")

add_text(slide, "Limitations actuelles",
         Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
         size=18, color=DANGER, bold=True)
add_bullets(slide, [
    "Scope académique : 4 notions × 5 modules",
    "Pas d'A/B test pédagogique réel",
    "Corpus RAG limité à 8 fichiers texte",
    "Pas d'agent générateur d'exercices",
    "Diagnostic adaptatif IRT-simplifié (pas complet)"
], Inches(0.5), Inches(2), Inches(6), Inches(4), size=13, bullet_color=DANGER)

add_text(slide, "Perspectives",
         Inches(7), Inches(1.3), Inches(6), Inches(0.5),
         size=18, color=SUCCESS, bold=True)
add_bullets(slide, [
    "Multi-langue (embeddings déjà multilingues)",
    "IRT complet (Item Response Theory)",
    "Intégration LMS (Moodle, Canvas) via MCP",
    "6e agent générateur d'exercices (LLM)",
    "Métriques d'efficacité pédagogique",
    "Déploiement cloud (Streamlit Cloud)"
], Inches(7), Inches(2), Inches(6), Inches(4), size=13, bullet_color=SUCCESS)

# ════════════════════════════════════════════════════════════════
# SLIDE 17 — Conclusion
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, PRIMARY)

add_text(slide, "🎓", Inches(5.5), Inches(0.8), Inches(2.5), Inches(1.5),
         size=80, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "Récapitulatif",
         Inches(1), Inches(2.3), Inches(11.3), Inches(0.7),
         size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

bullets_white = [
    "5 agents spécialisés orchestrés par LangGraph",
    "Conditional Edge démontrable (Sarah : -36% calculs)",
    "RAG vectoriel Chroma + embeddings multilingues",
    "3 outils MCP + Mind Layer LLM",
    "Communication A2A FIPA-ACL avec timeline visuelle",
    "15 tests pytest validés",
    "Conforme à tous les chapitres du cours"
]
tb = slide.shapes.add_textbox(Inches(2), Inches(3.2), Inches(9.3), Inches(3.5))
tf = tb.text_frame
tf.word_wrap = True
for i, item in enumerate(bullets_white):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(6)
    r1 = p.add_run()
    r1.text = "✓ "
    r1.font.size = Pt(16)
    r1.font.bold = True
    r1.font.color.rgb = ACCENT
    r2 = p.add_run()
    r2.text = item
    r2.font.size = Pt(16)
    r2.font.color.rgb = WHITE

# ════════════════════════════════════════════════════════════════
# SLIDE 18 — Questions
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)

add_text(slide, "?",
         Inches(5.5), Inches(1.5), Inches(2.3), Inches(2.5),
         size=180, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, "Questions ?",
         Inches(1), Inches(4), Inches(11.3), Inches(1),
         size=44, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, "Merci de votre attention",
         Inches(1), Inches(5.2), Inches(11.3), Inches(0.6),
         size=18, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

add_text(slide,
         "Équipe : Othmane · Dodo · Mohamed Yassir   |   Projet 4 · 4AISDR · Mai 2026",
         Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
         size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                            'Presentation_Projet4.pptx')
prs.save(output_path)

print(f"OK: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print(f"Size: {os.path.getsize(output_path) / 1024:.1f} KB")
